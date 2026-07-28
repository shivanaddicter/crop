import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_path):
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ── Color Palette (Dark Tech Theme) ──
    bg_color = RGBColor(15, 23, 42)        # Slate 900 (Dark Slate Blue/Grey)
    card_color = RGBColor(30, 41, 59)      # Slate 800 (Card background)
    border_color = RGBColor(71, 85, 105)   # Slate 600 (Card border)
    
    text_primary = RGBColor(255, 255, 255) # White
    text_secondary = RGBColor(203, 213, 225) # Slate 300 (Body text)
    text_muted = RGBColor(148, 163, 184)   # Slate 400 (Muted labels)
    accent_green = RGBColor(34, 197, 94)   # Emerald Green (Primary Accent)
    accent_green_light = RGBColor(187, 247, 208) # Mint Green (Secondary Accent)
    accent_blue = RGBColor(59, 130, 246)   # Blue Accent (for Tech stack)

    # ── Helper: Set Slide Background ──
    def set_slide_background(slide):
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = bg_color
        background.line.fill.background() # No border
        
        # Send to back is done implicitly as it's the first shape added
        # Add a subtle green glow circle at the bottom right
        glow = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(11.5), Inches(5.5), Inches(3.0), Inches(3.0)
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = RGBColor(20, 83, 45) # Dark green
        glow.line.fill.background()
        
    # ── Helper: Add Slide Header ──
    def add_header(slide, title_text, subtitle_text=None):
        # Title text box
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = text_primary
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = 'Segoe UI'
            p2.font.size = Pt(14)
            p2.font.color.rgb = accent_green
            p2.space_before = Pt(4)
            
        # Top divider line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.5), Inches(11.833), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = border_color
        line.line.fill.background()

    # ── Helper: Add Slide Footer ──
    def add_footer(slide):
        footer_box = slide.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(11.833), Inches(0.4))
        tf = footer_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = "Crop Yield Prediction AI  |  Hariprasath L"
        p.font.name = 'Segoe UI'
        p.font.size = Pt(10)
        p.font.color.rgb = text_muted
        p.alignment = PP_ALIGN.LEFT
        
        # Add slide number or small note on right
        # We can simulate this inside the same text frame or another

    # ── Helper: Create Card ──
    def add_card(slide, left, top, width, height, title, body_lines, icon=None, accent=accent_green):
        # Card shape
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = card_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        
        # Top highlight border strip
        top_strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.08)
        )
        top_strip.fill.solid()
        top_strip.fill.fore_color.rgb = accent
        top_strip.line.fill.background()
        
        # Content textbox
        tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(width - 0.5), Inches(height - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        # Title paragraph
        p_title = tf.paragraphs[0]
        title_text = f"{icon}  {title}" if icon else title
        p_title.text = title_text
        p_title.font.name = 'Trebuchet MS'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = text_primary
        p_title.space_after = Pt(12)
        
        # Body paragraphs
        for line in body_lines:
            p_body = tf.add_paragraph()
            p_body.text = line
            p_body.font.name = 'Segoe UI'
            p_body.font.size = Pt(12)
            p_body.font.color.rgb = text_secondary
            p_body.space_before = Pt(4)
            
        return card

    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)
    
    # Title Text Frame
    title_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "CROP YIELD PREDICTION AI"
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = text_primary
    p.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = "Intelligent Agronomic Recommendations & Real-Time Climate Analytics"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(22)
    p2.font.color.rgb = accent_green
    p2.space_after = Pt(40)
    
    # Subtitle line
    line = s1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.5), Inches(4.0), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent_green
    line.line.fill.background()
    
    # Author Info Box
    author_box = s1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(10.0), Inches(2.0))
    tf_auth = author_box.text_frame
    tf_auth.word_wrap = True
    
    p_auth = tf_auth.paragraphs[0]
    p_auth.text = "Developer & Founder:"
    p_auth.font.name = 'Segoe UI'
    p_auth.font.size = Pt(13)
    p_auth.font.color.rgb = text_muted
    p_auth.space_after = Pt(2)
    
    p_name = tf_auth.add_paragraph()
    p_name.text = "HARIPRASATH L"
    p_name.font.name = 'Trebuchet MS'
    p_name.font.size = Pt(18)
    p_name.font.bold = True
    p_name.font.color.rgb = text_primary
    
    p_inst = tf_auth.add_paragraph()
    p_inst.text = "B.Tech – Artificial Intelligence & Data Science (2nd Year)\nNPR College of Engineering and Technology, Madurai"
    p_inst.font.name = 'Segoe UI'
    p_inst.font.size = Pt(12)
    p_inst.font.color.rgb = text_secondary
    p_inst.space_before = Pt(4)

    # =========================================================================
    # SLIDE 2: The Agricultural Challenge (Problem Statement)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "The Agricultural Challenge", "Why modern farming requires a transition from guesswork to data science")
    add_footer(s2)
    
    # 3 Column layout
    card_width = 3.6
    card_height = 4.2
    top_pos = 2.1
    spacing = 0.5
    start_left = 0.75
    
    problems = [
        {
            "title": "Unpredictable Weather",
            "icon": "🌦️",
            "body": [
                "Climate change has broken traditional crop cycles.",
                "Unpredictable shifts in temperature, humidity, and erratic rainfall make intuition-based farming highly risky.",
                "Farmers struggle to adapt to localized weather shocks, leading to crop failures."
            ]
        },
        {
            "title": "Suboptimal Crop Selection",
            "icon": "❌",
            "body": [
                "Planting the wrong crop leads to severe yield losses.",
                "Crops require precise soil moisture and warmth thresholds to thrive.",
                "Without exact data analysis, resource utilization (water, fertilizers, labor) is highly inefficient."
            ]
        },
        {
            "title": "The Information Gap",
            "icon": "🔍",
            "body": [
                "Smallholder farmers lack access to localized, scientific data analysis.",
                "No direct tools exist to translate temperature and humidity values into immediate crop decisions.",
                "Bridges the information gap with simple, accessible web interfaces."
            ]
        }
    ]
    
    for i, prob in enumerate(problems):
        left_pos = start_left + i * (card_width + spacing)
        add_card(
            slide=s2,
            left=left_pos,
            top=top_pos,
            width=card_width,
            height=card_height,
            title=prob["title"],
            body_lines=prob["body"],
            icon=prob["icon"]
        )

    # =========================================================================
    # SLIDE 3: The Smart Farming Solution
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "The Smart Farming Solution", "Introducing Crop Yield Prediction AI — a scientific approach to agriculture")
    add_footer(s3)
    
    # Left Hero Card
    hero_card = s3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.1), Inches(5.2), Inches(4.2)
    )
    hero_card.fill.solid()
    hero_card.fill.fore_color.rgb = card_color
    hero_card.line.color.rgb = accent_green
    hero_card.line.width = Pt(1.5)
    
    tb_hero = s3.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(4.7), Inches(3.6))
    tf_hero = tb_hero.text_frame
    tf_hero.word_wrap = True
    
    p_hero_title = tf_hero.paragraphs[0]
    p_hero_title.text = "🎯 Project Objective"
    p_hero_title.font.name = 'Trebuchet MS'
    p_hero_title.font.size = Pt(20)
    p_hero_title.font.bold = True
    p_hero_title.font.color.rgb = text_primary
    p_hero_title.space_after = Pt(14)
    
    p_hero_body = tf_hero.add_paragraph()
    p_hero_body.text = "To build a robust, responsive web application that leverages weather parameters to output exact crop recommendations, minimizing the risk of failure while improving resource planning."
    p_hero_body.font.name = 'Segoe UI'
    p_hero_body.font.size = Pt(14)
    p_hero_body.font.color.rgb = text_secondary
    p_hero_body.space_after = Pt(14)
    
    p_hero_sub = tf_hero.add_paragraph()
    p_hero_sub.text = "By collecting user data (Temperature, Rainfall, Humidity) or fetching real-time data dynamically, the system simulates agronomic expert models to recommend Rice, Wheat, Sugarcane, or Maize alongside comprehensive care guidelines."
    p_hero_sub.font.name = 'Segoe UI'
    p_hero_sub.font.size = Pt(13)
    p_hero_sub.font.color.rgb = text_muted

    # Right side 3 features list
    right_start_left = 6.4
    card_w = 6.1
    card_h = 1.25
    spacing_h = 0.225
    
    features = [
        ("Precision Agriculture", "Uses rule-based ML modeling to calculate crop compatibility using strict environmental parameters.", "🌾"),
        ("Real-time Local Adaptability", "Fetches actual weather data via OpenWeatherMap API for dynamic, localized calculations.", "☁️"),
        ("Scalable Farm Management", "Allows CSV file batch uploads to process recommendations for multiple acreage plots at once.", "📈")
    ]
    
    for i, (title, desc, icon) in enumerate(features):
        top = 2.1 + i * (card_h + spacing_h)
        add_card(
            slide=s3,
            left=right_start_left,
            top=top,
            width=card_w,
            height=card_h,
            title=title,
            body_lines=[desc],
            icon=icon
        )

    # =========================================================================
    # SLIDE 4: Key Platform Features & Capabilities
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "Key Platform Capabilities", "A modern toolset designed for farmers and agronomy planners")
    add_footer(s4)
    
    # 2x2 Grid Layout
    w_card = 5.6
    h_card = 2.0
    top_pos_r1 = 2.1
    top_pos_r2 = 4.35
    left_c1 = 0.75
    left_c2 = 6.98
    
    add_card(
        slide=s4, left=left_c1, top=top_pos_r1, width=w_card, height=h_card,
        title="Predictive Crop Model",
        body_lines=["Input Temperature, Rainfall, and Humidity manually.", "Instantly outputs the recommended crop.", "Provides dedicated recommendations for sowing, watering, and harvesting."],
        icon="🔮"
    )
    
    add_card(
        slide=s4, left=left_c2, top=top_pos_r1, width=w_card, height=h_card,
        title="Live Weather Predictor",
        body_lines=["Type any city name (e.g., Chennai, Madurai).", "Integrates external APIs to fetch actual weather indicators.", "Performs instant prediction on real-world, real-time metrics."],
        icon="⛅"
    )
    
    add_card(
        slide=s4, left=left_c1, top=top_pos_r2, width=w_card, height=h_card,
        title="Batch Processing (CSV Upload)",
        body_lines=["Upload a CSV file containing multiple land plot entries.", "Processes all records in a single click.", "Includes validation to automatically log successful and skipped rows."],
        icon="📁"
    )
    
    add_card(
        slide=s4, left=left_c2, top=top_pos_r2, width=w_card, height=h_card,
        title="Interactive Analytics Dashboard",
        body_lines=["Aggregates predictions into simple average metrics.", "Features dynamic Chart.js charts showing temperature/rainfall trends.", "Integrates doughnut charts showing crop recommendation distributions."],
        icon="📊"
    )

    # =========================================================================
    # SLIDE 5: System Architecture & Technology Stack
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Architecture & Tech Stack", "Modern, scalable, and responsive full-stack implementation")
    add_footer(s5)
    
    # Left Box - Backend
    back_box = s5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.1), Inches(3.6), Inches(4.2)
    )
    back_box.fill.solid()
    back_box.fill.fore_color.rgb = card_color
    back_box.line.color.rgb = border_color
    
    tb_b = s5.shapes.add_textbox(Inches(0.95), Inches(2.3), Inches(3.2), Inches(3.8))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    p_bt = tf_b.paragraphs[0]
    p_bt.text = "⚡ Backend Server"
    p_bt.font.name = 'Trebuchet MS'
    p_bt.font.size = Pt(18)
    p_bt.font.bold = True
    p_bt.font.color.rgb = accent_green
    p_bt.space_after = Pt(12)
    
    backend_bullets = [
        "Python (Core Logic)",
        "Django Framework",
        "Django ORM (Object Relational Mapping)",
        "Built-in Authentication Session Guarding",
        "SQLite Database (Quick, light storage)",
        "Django Messages Framework"
    ]
    for bullet in backend_bullets:
        pb = tf_b.add_paragraph()
        pb.text = f"• {bullet}"
        pb.font.name = 'Segoe UI'
        pb.font.size = Pt(12)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(4)
        
    # Middle Box - Frontend
    front_box = s5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.85), Inches(2.1), Inches(3.6), Inches(4.2)
    )
    front_box.fill.solid()
    front_box.fill.fore_color.rgb = card_color
    front_box.line.color.rgb = border_color
    
    tb_fr = s5.shapes.add_textbox(Inches(5.05), Inches(2.3), Inches(3.2), Inches(3.8))
    tf_fr = tb_fr.text_frame
    tf_fr.word_wrap = True
    p_ft = tf_fr.paragraphs[0]
    p_ft.text = "🎨 Presentation Layer"
    p_ft.font.name = 'Trebuchet MS'
    p_ft.font.size = Pt(18)
    p_ft.font.bold = True
    p_ft.font.color.rgb = RGBColor(59, 130, 246) # Blue accent
    p_ft.space_after = Pt(12)
    
    frontend_bullets = [
        "Semantic HTML5 (Clean layout structure)",
        "CSS3 Custom Variables (Dynamic styling)",
        "Custom Dark Cyber Theme",
        "Chart.js (Interactive UI data visualizations)",
        "Bootstrap (Responsive navigation grids)",
        "FontAwesome / Custom SVG icons"
    ]
    for bullet in frontend_bullets:
        pb = tf_fr.add_paragraph()
        pb.text = f"• {bullet}"
        pb.font.name = 'Segoe UI'
        pb.font.size = Pt(12)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(4)
        
    # Right Box - Core Engine & Services
    eng_box = s5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.95), Inches(2.1), Inches(3.6), Inches(4.2)
    )
    eng_box.fill.solid()
    eng_box.fill.fore_color.rgb = card_color
    eng_box.line.color.rgb = border_color
    
    tb_e = s5.shapes.add_textbox(Inches(9.15), Inches(2.3), Inches(3.2), Inches(3.8))
    tf_e = tb_e.text_frame
    tf_e.word_wrap = True
    p_et = tf_e.paragraphs[0]
    p_et.text = "⚙️ Integration & Logic"
    p_et.font.name = 'Trebuchet MS'
    p_et.font.size = Pt(18)
    p_et.font.bold = True
    p_et.font.color.rgb = RGBColor(234, 179, 8) # Yellow Accent
    p_et.space_after = Pt(12)
    
    engine_bullets = [
        "Rule-Based ML Crop Decision Tree",
        "OpenWeatherMap API Connection",
        "JSON-based weather schema conversion",
        "CSV DictReader upload parser",
        "CSV writer history exporter",
        "Python Requests module"
    ]
    for bullet in engine_bullets:
        pb = tf_e.add_paragraph()
        pb.text = f"• {bullet}"
        pb.font.name = 'Segoe UI'
        pb.font.size = Pt(12)
        pb.font.color.rgb = text_secondary
        pb.space_before = Pt(4)

    # =========================================================================
    # SLIDE 6: Under the Hood: Agronomic Rule Engine
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "Agronomic Rule-Based ML Model", "How the environmental model matches climate metrics to specific crops")
    add_footer(s6)
    
    card_w = 2.7
    card_h = 4.2
    top = 2.1
    spacing = 0.3
    start_left = 0.75
    
    crops = [
        {
            "name": "🌾 RICE",
            "color": accent_green,
            "body": [
                "Water Intensive",
                "• Rain > 200mm (when Temp >= 20°C)",
                "• Moderate Temp (20-30°C) with high Humidity (>75%)",
                "• Advisory: Ideal for clayey soils that retain moisture."
            ]
        },
        {
            "name": "🌾 WHEAT",
            "color": RGBColor(59, 130, 246), # Blue
            "body": [
                "Cool / Dry Tolerant",
                "• Rain > 200mm (when Temp < 20°C)",
                "• Temp < 20°C in general",
                "• Moderate Temp (20-30°C) with low Rain & Humidity",
                "• Advisory: Thrives in loamy, well-drained soils."
            ]
        },
        {
            "name": "🍬 SUGARCANE",
            "color": RGBColor(236, 72, 153), # Pink
            "body": [
                "Heat / Arid Tolerant",
                "• High Temp > 32°C",
                "• Low Humidity < 40%",
                "• Advisory: Requires deep, fertile soils and timely irrigation cycles."
            ]
        },
        {
            "name": "🌽 MAIZE",
            "color": RGBColor(234, 179, 8), # Yellow
            "body": [
                "Moderate Climates",
                "• High Temp > 32°C with Humidity >= 40%",
                "• Moderate Temp (20-30°C) with Rain > 100mm",
                "• Advisory: Requires good drainage and rich nitrogen levels."
            ]
        }
    ]
    
    for i, crop in enumerate(crops):
        left_pos = start_left + i * (card_w + spacing)
        add_card(
            slide=s6,
            left=left_pos,
            top=top,
            width=card_w,
            height=card_h,
            title=crop["name"],
            body_lines=crop["body"],
            accent=crop["color"]
        )

    # =========================================================================
    # SLIDE 7: Platform Features & Operations
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "Database Design & Data Flow", "How data flows through the SQLite store and CSV channels")
    add_footer(s7)
    
    # Left Card: Database Models
    add_card(
        slide=s7, left=0.75, top=2.1, width=5.6, height=4.2,
        title="Database Models (Django Models)",
        body_lines=[
            "📁 CropData Table:",
            "  • user (ForeignKey to User model)",
            "  • temperature (Float)",
            "  • rainfall (Float)",
            "  • humidity (Float)",
            "  • result (String - e.g. 'Rice')",
            "  • created_at (DateTimeField - auto timestamp)",
            "",
            "📁 UploadedCSV Table:",
            "  • user (ForeignKey to User)",
            "  • file (FileField - path to media folder)",
            "  • uploaded_at (DateTimeField)"
        ],
        icon="📂"
    )
    
    # Right Card: Data Operations
    add_card(
        slide=s7, left=6.98, top=2.1, width=5.6, height=4.2,
        title="Core Data Workflows",
        body_lines=[
            "🔐 Authentication Flow:",
            "  • Registration (UserCreationForm) enforces security.",
            "  • Route-level @login_required decorator restricts analytical access.",
            "",
            "📥 CSV Exporting Engine:",
            "  • Queries all CropData records for logged-in user.",
            "  • Emits streaming CSV response object with headers in real time.",
            "",
            "📤 CSV Importing Engine:",
            "  • Resets file buffer and parses using DictReader.",
            "  • Iterates, converts strings to floats, handles missing entries.",
            "  • Feeds values to predictor and commits batch rows."
        ],
        icon="🔄"
    )

    # =========================================================================
    # SLIDE 8: Developer & Founder Profile
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Meet the Founder & Developer", "Designing technology solutions for real-world impact")
    add_footer(s8)
    
    # Profile Card (Left)
    prof_card = s8.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.1), Inches(5.6), Inches(4.2)
    )
    prof_card.fill.solid()
    prof_card.fill.fore_color.rgb = card_color
    prof_card.line.color.rgb = border_color
    
    tb_p = s8.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.1), Inches(3.6))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True
    
    p_name = tf_p.paragraphs[0]
    p_name.text = "👤 HARIPRASATH L"
    p_name.font.name = 'Trebuchet MS'
    p_name.font.size = Pt(22)
    p_name.font.bold = True
    p_name.font.color.rgb = text_primary
    
    p_title = tf_p.add_paragraph()
    p_title.text = "Full Stack Web Developer & AI/DS Student"
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = accent_green
    p_title.space_after = Pt(14)
    
    p_bio = tf_p.add_paragraph()
    p_bio.text = "Enthusiastic and driven 2nd year B.Tech student specializing in Artificial Intelligence and Data Science at NPR College of Engineering and Technology, Madurai.\nHands-on experience building full stack web projects and exploring agricultural analytics using Python.\nActively participates in hackathons, engineering symposia, and AI workshops."
    p_bio.font.name = 'Segoe UI'
    p_bio.font.size = Pt(12)
    p_bio.font.color.rgb = text_secondary
    p_bio.space_before = Pt(6)

    # Skills Card (Right)
    skills_card = s8.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(2.1), Inches(5.6), Inches(4.2)
    )
    skills_card.fill.solid()
    skills_card.fill.fore_color.rgb = card_color
    skills_card.line.color.rgb = border_color
    
    tb_s = s8.shapes.add_textbox(Inches(7.23), Inches(2.4), Inches(5.1), Inches(3.6))
    tf_s = tb_s.text_frame
    tf_s.word_wrap = True
    
    p_sk_t = tf_s.paragraphs[0]
    p_sk_t.text = "🛠️ Technical Core & Contact"
    p_sk_t.font.name = 'Trebuchet MS'
    p_sk_t.font.size = Pt(18)
    p_sk_t.font.bold = True
    p_sk_t.font.color.rgb = text_primary
    p_sk_t.space_after = Pt(12)
    
    skills_info = [
        "Programming: Python, Java, SQL",
        "Frontend: HTML5, CSS3 Custom styling, JavaScript, Bootstrap, jQuery",
        "Backend & Web: Flask, Django frameworks, Django ORM",
        "Database: MySQL, MongoDB, SQLite",
        "Location: Madurai, Tamil Nadu, India",
        "Contact: hariprasath72788@gmail.com",
        "Projects: github.com/shivanaddicter | linkedin.com/in/hariprasath-l-5b6b40312"
    ]
    for skill in skills_info:
        ps = tf_s.add_paragraph()
        ps.text = f"• {skill}"
        ps.font.name = 'Segoe UI'
        ps.font.size = Pt(12)
        ps.font.color.rgb = text_secondary
        ps.space_before = Pt(4)

    # =========================================================================
    # SLIDE 9: Future Roadmap & Future Scope
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "Future Roadmap & Project Scaling", "Key strategies to transition from rule-based simulations to automated systems")
    add_footer(s9)
    
    # 3 Column layout
    card_width = 3.6
    card_height = 4.2
    top_pos = 2.1
    spacing = 0.5
    start_left = 0.75
    
    roadmap = [
        {
            "title": "Machine Learning Models",
            "icon": "🧠",
            "body": [
                "Transition from static rule-based boundaries to active ML.",
                "Integrate scikit-learn models like Random Forest, SVM, or XGBoost.",
                "Train on extensive agricultural datasets containing multiple features (soil NPK, soil pH, organic carbon levels)."
            ]
        },
        {
            "title": "Automated Weather Services",
            "icon": "📍",
            "body": [
                "Incorporate geolocation in the client-side browser.",
                "Fetch the exact latitude and longitude coordinates automatically.",
                "Call the weather API dynamically to produce recommendation values without typing in city names manually."
            ]
        },
        {
            "title": "Mobile App Deployment",
            "icon": "📱",
            "body": [
                "Develop a mobile version of the crop yield prediction portal.",
                "Provide lightweight offline support to work under low connectivity.",
                "Incorporate regional Indian languages (e.g. Tamil) to maximize accessibility for local farmers."
            ]
        }
    ]
    
    for i, road in enumerate(roadmap):
        left_pos = start_left + i * (card_width + spacing)
        add_card(
            slide=s9,
            left=left_pos,
            top=top_pos,
            width=card_width,
            height=card_height,
            title=road["title"],
            body_lines=road["body"],
            icon=road["icon"]
        )
        
    # Save presentation
    prs.save(output_path)
    print(f"Presentation created successfully at: {output_path}")

if __name__ == "__main__":
    out_dir = r"c:\Users\Lenovo\Desktop"
    if len(sys.argv) > 1:
        out_dir = sys.argv[1]
    
    filename = "Crop_Yield_Prediction_AI_Presentation.pptx"
    output_filepath = os.path.join(out_dir, filename)
    create_presentation(output_filepath)
